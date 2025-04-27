"""Action to fetch Confluence pages and create PowerPoint slide decks."""

import os
import re
import json
import tempfile
from pathlib import Path
from dotenv import load_dotenv
from typing import Dict, List, Tuple, Any

from solace_ai_connector.common.log import log
from solace_agent_mesh.common.action import Action
from solace_agent_mesh.common.action_response import ActionResponse

# Import functions from plm_writing_assistant
from plm_writing_assistant.actions.plm_writing_assistant import PlmWritingAssistant

# Import pptx for PowerPoint creation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load environment variables
load_dotenv()


class PLMSlideDeckBuilder(Action):
    def __init__(self, **kwargs):
        super().__init__(
            {
                "name": "plm_slide_deck_builder",
                "description": "Fetches Confluence pages and creates PowerPoint slide decks.",
                "prompt_directive": (
                    "Provide a Confluence page URL to fetch its content and generate a PowerPoint presentation."
                ),
                "params": [
                    {
                        "name": "confluence_url",
                        "desc": "URL of the Confluence page to format into slides",
                        "type": "string",
                    },
                    {
                        "name": "slides_count",
                        "desc": "The target number of slides to create",
                        "type": "int",
                        "default": 10
                    },
                    {
                        "name": "output_path",
                        "desc": "Path where the PowerPoint file should be saved",
                        "type": "string",
                        "default": "presentation.pptx"
                    },
                    {
                        "name": "design_style",
                        "desc": "Style hint for the presentation design (e.g., 'corporate', 'creative', 'technical')",
                        "type": "string",
                        "default": "professional"
                    }
                ],
                "required_scopes": ["plm_slide_deck_builder:create:write"],
            },
            **kwargs,
        )
        # Create an instance of PlmWritingAssistant to reuse its methods
        self.writing_assistant = PlmWritingAssistant(**kwargs)

    def invoke(self, params, meta={}) -> ActionResponse:
        # Get parameters
        confluence_url = params.get("confluence_url")
        slides_count = params.get("slides_count", 10)
        output_path = params.get("output_path", "presentation.pptx")
        design_style = params.get("design_style", "professional")
        
        if not confluence_url:
            return ActionResponse(message="Error: No Confluence URL provided")
        
        # Fetch the Confluence page content using the writing assistant's method
        log.info("Fetching Confluence page: %s", confluence_url)
        content = self.writing_assistant.fetch_confluence_page(confluence_url)
        
        # If fetching failed, return the error message
        if content.startswith("⚠️"):
            return ActionResponse(message=content)
        
        # Ensure API token is properly formatted (strip quotes if present)
        api_token = os.getenv("ATLASSIAN_API_TOKEN", "")
        if api_token and (api_token.startswith('"') and api_token.endswith('"')):
            api_token = api_token[1:-1]  # Strip surrounding quotes
            os.environ["ATLASSIAN_API_TOKEN"] = api_token
            log.info("Stripped quotes from ATLASSIAN_API_TOKEN")
        
        # Format content for slides
        formatted_slides = self.format_for_slides(content, slides_count)
        
        # Log Confluence credentials for debugging
        log.info("Using Confluence credentials:")
        log.info("CONFLUENCE_BASE_URL: %s", os.getenv("CONFLUENCE_BASE_URL"))
        log.info("ATLASSIAN_EMAIL: %s", os.getenv("ATLASSIAN_EMAIL"))
        log.info("ATLASSIAN_API_TOKEN length: %s", len(os.getenv("ATLASSIAN_API_TOKEN", "")))
        
        # Determine design choices using LLM
        design_choices = self.determine_design_choices(content, design_style)
        
        # Create PowerPoint presentation
        ppt_path = self.create_powerpoint(formatted_slides, design_choices, output_path)
        
        return ActionResponse(
            message=f"PowerPoint presentation created successfully at {ppt_path}",
            data={
                "presentation_path": ppt_path,
                "slide_count": slides_count,
                "design_style": design_style
            }
        )

    def format_for_slides(self, content, slides_count):
        """Format content for slide deck presentation using the LLM service"""
        if self.agent and hasattr(self.agent, "do_llm_service_request"):
            messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a presentation expert. Your task is to transform Confluence content into "
                        "well-structured slide content that can be easily parsed for a PowerPoint presentation. "
                        "Focus on creating professional, concise slides with clear messaging."
                    )
                },
                {
                    "role": "user",
                    "content": (
                        f"Create a complete {slides_count}-slide presentation based on the following Confluence content. "
                        f"Use Solace branding colors (#00C895 green, #2D3142 dark blue/gray, #00A6ED blue).\n\n"
                        f"For each slide, provide:\n"
                        f"1. A clear, concise title that captures the main point\n"
                        f"2. 3-5 bullet points with key information (keep each point under 15 words)\n"
                        f"3. Visual element suggestions marked as [VISUAL: description]\n\n"
                        f"Start with a title slide, then an agenda/overview slide, followed by content slides, "
                        f"and end with a summary/next steps slide.\n\n"
                        f"Format each slide as:\n\n"
                        f"## Slide X: Title\n"
                        f"- Key point 1\n"
                        f"- Key point 2\n"
                        f"- Key point 3\n"
                        f"[VISUAL: Brief description of recommended visual]\n\n"
                        f"Here's the content to transform:\n\n{content}"
                    )
                }
            ]
            
            try:
                log.info("Sending content to LLM service for slide formatting")
                response = self.agent.do_llm_service_request(messages)
                formatted_content = response.get("content", "")
                log.info("Successfully formatted content for slides")
                return formatted_content
            except Exception as e:
                error_message = f"Error calling LLM service: {str(e)}"
                log.error(error_message)
                return error_message
        else:
            return "Error: LLM service not available"
            
    def determine_design_choices(self, content, design_style):
        """Use LLM to determine design choices for the presentation"""
        # If design_style contains 'solace', use Solace branding
        if "solace" in design_style.lower():
            log.info("Using Solace branding colors")
            return {
                "primary_color": "#00C895",  # Solace green
                "secondary_color": "#2D3142", # Dark blue/gray
                "accent_color": "#00A6ED",    # Solace blue
                "font_title": "Arial",
                "font_body": "Arial",
                "title_size": 28,
                "body_size": 18,
                "layout_style": "corporate",
                "background_style": "white with subtle green accents"
            }
            
        if self.agent and hasattr(self.agent, "do_llm_service_request"):
            messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a presentation design expert. Your task is to determine the best design choices "
                        "for a PowerPoint presentation based on the content and requested style. "
                        "If the style mentions 'solace' or 'solace.com', use Solace's brand colors: "
                        "primary: #00C895 (green), secondary: #2D3142 (dark blue/gray), accent: #00A6ED (blue)."
                    )
                },
                {
                    "role": "user",
                    "content": (
                        f"Based on this content excerpt and the requested design style '{design_style}', "
                        f"provide design recommendations for a PowerPoint presentation.\n\n"
                        f"Content excerpt (first 500 chars):\n{content[:500]}...\n\n"
                        f"Return your recommendations as a JSON object with these fields:\n"
                        f"- primary_color: RGB hex code (e.g., '#0070C0')\n"
                        f"- secondary_color: RGB hex code\n"
                        f"- accent_color: RGB hex code\n"
                        f"- font_title: font name for titles\n"
                        f"- font_body: font name for body text\n"
                        f"- title_size: point size for titles (e.g., 28)\n"
                        f"- body_size: point size for body text (e.g., 18)\n"
                        f"- layout_style: 'minimal', 'corporate', or 'creative'\n"
                        f"- background_style: description of background style"
                    )
                }
            ]
            
            try:
                log.info("Sending request to LLM service for design choices")
                response = self.agent.do_llm_service_request(messages)
                design_content = response.get("content", "")
                
                # Extract JSON from the response
                json_match = re.search(r'```json\s*(.*?)\s*```', design_content, re.DOTALL)
                if json_match:
                    design_json = json_match.group(1)
                else:
                    design_json = design_content
                
                # Parse the JSON
                try:
                    design_choices = json.loads(design_json)
                    log.info("Successfully determined design choices")
                    return design_choices
                except json.JSONDecodeError:
                    log.error("Failed to parse design choices JSON")
                    # Return default design choices
                    return {
                        "primary_color": "#00C895",  # Solace green
                        "secondary_color": "#2D3142", # Dark blue/gray
                        "accent_color": "#00A6ED",    # Solace blue
                        "font_title": "Calibri",
                        "font_body": "Calibri",
                        "title_size": 28,
                        "body_size": 18,
                        "layout_style": "corporate",
                        "background_style": "solid white"
                    }
            except Exception as e:
                error_message = f"Error calling LLM service for design: {str(e)}"
                log.error(error_message)
                # Return default design choices
                return {
                    "primary_color": "#00C895",  # Solace green
                    "secondary_color": "#2D3142", # Dark blue/gray
                    "accent_color": "#00A6ED",    # Solace blue
                    "font_title": "Calibri",
                    "font_body": "Calibri",
                    "title_size": 28,
                    "body_size": 18,
                    "layout_style": "corporate",
                    "background_style": "solid white"
                }
        else:
            log.error("LLM service not available for design choices")
            # Return default design choices
            return {
                "primary_color": "#00C895",  # Solace green
                "secondary_color": "#2D3142", # Dark blue/gray
                "accent_color": "#00A6ED",    # Solace blue
                "font_title": "Calibri",
                "font_body": "Calibri",
                "title_size": 28,
                "body_size": 18,
                "layout_style": "corporate",
                "background_style": "solid white"
            }
            
    def create_powerpoint(self, formatted_slides, design_choices, output_path):
        """Create a PowerPoint presentation from the formatted slides and design choices"""
        try:
            # Log the formatted slide content for debugging
            log.info("Formatted slide content generated by LLM:\n%s", formatted_slides[:500] + "...")
            
            # Create a new presentation
            prs = Presentation()
            
            # Parse the formatted slides
            slides = self.parse_formatted_slides(formatted_slides)
            
            # Extract design choices
            primary_color = self.hex_to_rgb(design_choices.get("primary_color", "#0070C0"))
            secondary_color = self.hex_to_rgb(design_choices.get("secondary_color", "#4472C4"))
            accent_color = self.hex_to_rgb(design_choices.get("accent_color", "#70AD47"))
            font_title = design_choices.get("font_title", "Calibri")
            font_body = design_choices.get("font_body", "Calibri")
            title_size = design_choices.get("title_size", 28)
            body_size = design_choices.get("body_size", 18)
            
            # Create title slide
            if slides:
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                # Get the first slide's title as the presentation title
                first_slide = slides[0]
                title.text = first_slide["title"].replace("Slide 1: ", "")
                
                # Set a subtitle
                subtitle.text = "Generated from Confluence"
                
                # Apply design to title slide
                title.text_frame.paragraphs[0].font.name = font_title
                title.text_frame.paragraphs[0].font.size = Pt(title_size + 8)
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_color)
                
                subtitle.text_frame.paragraphs[0].font.name = font_body
                subtitle.text_frame.paragraphs[0].font.size = Pt(body_size + 2)
                subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(*secondary_color)
            
            # Create content slides
            for i, slide_content in enumerate(slides):
                # Skip the first slide as we used it for the title
                if i == 0:
                    continue
                    
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = slide_content["title"]
                
                # Apply design to title
                title.text_frame.paragraphs[0].font.name = font_title
                title.text_frame.paragraphs[0].font.size = Pt(title_size)
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_color)
                
                # Add content as bullet points
                content = slide.placeholders[1]
                tf = content.text_frame
                
                for bullet_point in slide_content["bullets"]:
                    p = tf.add_paragraph()
                    p.text = bullet_point
                    p.level = 0
                    p.font.name = font_body
                    p.font.size = Pt(body_size)
                    p.font.color.rgb = RGBColor(*secondary_color)
                
                # If there's a visual suggestion, add it as a note
                if slide_content.get("visual"):
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = f"Visual suggestion: {slide_content['visual']}"
            
            # Save the presentation
            prs.save(output_path)
            log.info(f"PowerPoint presentation saved to {output_path}")
            return output_path
            
        except Exception as e:
            error_message = f"Error creating PowerPoint: {str(e)}"
            log.error(error_message)
            return f"Failed to create presentation: {error_message}"
    
    def parse_formatted_slides(self, formatted_text):
        """Parse the formatted slide text into a structured format"""
        slides = []
        current_slide = None
        
        # Split the text by lines
        lines = formatted_text.strip().split('\n')
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
                
            # Check if this is a slide title
            slide_match = re.match(r'^##\s+Slide\s+(\d+):\s+(.*)', line)
            if slide_match:
                # If we already have a slide, add it to the list
                if current_slide:
                    slides.append(current_slide)
                
                # Start a new slide
                current_slide = {
                    "title": line.replace("## ", ""),
                    "bullets": [],
                    "visual": None
                }
                continue
                
            # Check if this is a bullet point
            if line.startswith('- '):
                if current_slide:
                    current_slide["bullets"].append(line[2:])
                continue
                
            # Check if this is a visual suggestion
            visual_match = re.match(r'\[VISUAL:\s+(.*)\]', line)
            if visual_match and current_slide:
                current_slide["visual"] = visual_match.group(1)
                continue
        
        # Add the last slide if there is one
        if current_slide:
            slides.append(current_slide)
            
        return slides
    
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        # Remove the # if present
        hex_color = hex_color.lstrip('#')
        
        # Convert to RGB
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return (r, g, b)
        except (ValueError, IndexError):
            # Return a default blue color if conversion fails
            return (0, 112, 192)  # Default blue
